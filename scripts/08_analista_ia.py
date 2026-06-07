import argparse
import base64
import csv
import glob
import hashlib
import json
import os
import re
import sys
import time
import subprocess
from datetime import datetime

import requests

# ==========================================
# CONFIGURACIÓN DEL LLM (Valores default — sobreescribibles por args)
# ==========================================
OLLAMA_BASE_URL    = "http://localhost:11434"   # Host Ollama (local o remoto)
OLLAMA_URL         = OLLAMA_BASE_URL + "/api/generate"
MODELO_LLM         = "gemma3:4b"                # Modelo por defecto (RPi5)
DIRECTORIO_DEFAULT = "/mnt/Destino_ForenSys"

# --- PERFIL: Raspberry Pi 5 (8 GB RAM) ---
RPI_CTX     = 2048   # Ventana de contexto conservadora
RPI_THREAD  = 4      # 4 cores físicos
RPI_TIMEOUT = 900    # 15 minutos

# --- PERFIL: PC Escritorio (Ryzen 5600G + RX 6600 XT) ---
PC_CTX     = 4096   # Ventana de contexto amplia
PC_THREAD  = 12     # 6c/12t del 5600G
PC_TIMEOUT = 300    # 5 minutos (GPU/CPU más rápida)

# Valores activos (se sobrescriben en main según el motor elegido)
NUM_CTX    = RPI_CTX
NUM_THREAD = RPI_THREAD
KEEP_ALIVE = "0"     # Descargar modelo tras responder
TIMEOUT_SOLICITUD = RPI_TIMEOUT

# Presupuesto de caracteres del prompt
MAX_CARACTERES_EVIDENCIA = 6000

# ==========================================
# WHITELIST DE MODELOS AUDITADOS (CRÍTICA 5)
# Solo estos modelos están validados para uso forense.
# ==========================================
MODELOS_AUDITADOS = {
    'gemma3:4b',
    'gemma3:1b',
    'gemma:4b',
    'gemma2:4b',
    'gemma4:4b',
    'gemma4:e4b',
    'gemma4-forense:latest',
    'gemma4-e4b-uncensored:latest',
    'gemma:2b',
    'llama3.2:3b',
    'llama3.2:1b',
    'mistral:7b',
    'mistral:latest',
    'llama3:8b',
    'llama3:latest',
    'phi3.5',
    'phi3.5:latest',
}

# ==========================================
# DISCLAIMER LEGAL — CRÍTICA 3
# Visible al inicio de cada síntesis generada.
# ==========================================
DISCLAIMER_LEGAL = """
╔══════════════════════════════════════════════════════════════════════════════╗
║              ⚠️  AVISO LEGAL — DOCUMENTO DE ASISTENCIA ⚠️                  ║
╠══════════════════════════════════════════════════════════════════════════════╣
║                                                                              ║
║  ESTE DOCUMENTO ES UNA SÍNTESIS AUTOMÁTICA GENERADA POR INTELIGENCIA        ║
║  ARTIFICIAL. NO ES UN DICTAMEN PERICIAL CERTIFICADO.                        ║
║                                                                              ║
║  LIMITACIONES CRÍTICAS:                                                      ║
║  • Basado en algoritmos de machine learning (LLM). Puede contener errores.  ║
║  • Puede generar "alucinaciones" — afirmaciones plausibles pero falsas.      ║
║  • NO es admisible como evidencia directa en procedimientos legales.         ║
║  • REQUIERE validación y firma de un perito forense certificado.             ║
║                                                                              ║
║  USO PERMITIDO:                                                              ║
║  ✓ Apoyo en investigación preliminar                                         ║
║  ✓ Síntesis y organización de múltiples fuentes de datos                    ║
║  ✓ Identificación de áreas que requieren análisis forense profundo           ║
║                                                                              ║
║  USO PROHIBIDO:                                                              ║
║  ✗ Dictamen pericial sin revisión y firma de perito                         ║
║  ✗ Evidencia directa en procedimientos judiciales                           ║
║  ✗ Conclusiones definitivas sin análisis técnico independiente               ║
║                                                                              ║
║  RESPONSABILIDAD: El operador es responsable de verificar cada hallazgo.    ║
║                                                                              ║
╚══════════════════════════════════════════════════════════════════════════════╝
"""

# Prompt del asistente — rol honesto, sin fingir ser perito (CRÍTICA 1 y 8)
PROMPT_ASISTENTE = """Eres un ASISTENTE DE ANÁLISIS FORENSE (no un perito certificado).

Tu rol es sintetizar y organizar la evidencia digital para facilitar el trabajo \
del perito humano. NO emites conclusiones legales ni reemplazas el análisis de un perito.

Cada hallazgo debe indicar:
1. Fuente de datos exacta
2. Nivel de confianza: BAJO / MEDIO / ALTO
3. Si requiere validación manual por el perito

EVIDENCIA RECOPILADA:
{evidencia_cruda}

INSTRUCCIONES (máximo 1000 palabras, formato Markdown):

## 1. Identificación del Equipo
- Nombre del equipo, zona horaria, sistema operativo detectado.
- Dispositivos USB que fueron conectados (marca, serial).

## 2. Usuarios del Sistema
- Lista de usuarios detectados con su nivel de actividad.
- Documentos recientes, comandos ejecutados, programas usados por cada uno.

## 3. Perfil de Uso
- Descripción del perfil del usuario principal basándose en: programas instalados, \
historial web, tipos de archivos. (Nivel de confianza a indicar)

## 4. Cronología de Actividad Relevante
- Las 5-10 acciones más relevantes en orden cronológico.
- Incluye: URLs visitadas, programas ejecutados, archivos descargados.

## 5. Alertas de Seguridad
- Archivos con alta entropía (posible ransomware/cifrado).
- Mecanismos de persistencia sospechosos (auto-arranque fuera de rutas confiables).
- Extensiones falsas o técnicas antiforenses detectadas.

## 6. Evidencia Multimedia
- Archivos con coordenadas GPS embebidas.
- Cámaras/dispositivos detectados en metadatos EXIF.

## 7. Síntesis Final (NO es conclusión pericial)
- Párrafo final resumiendo hallazgos y señalando qué áreas requieren \
investigación adicional por parte del perito. Incluir nivel de confianza global.

REGLAS CRÍTICAS DE SEGURIDAD (GRADO 0.0 ALUCINACIONES):
- NO inventes ni asumas datos, horas, archivos o usuarios. Solo analiza la evidencia proporcionada.
- BAJO NINGUNA CIRCUNSTANCIA alteres o cambies la información técnica (hashes, IPs, rutas, fechas).
- Si no hay información suficiente sobre un punto, indicar explícitamente "Sin datos disponibles".
- Manten un tono estrictamente formal y técnico. Indicar la fuente exacta de cada afirmación.
- Esta síntesis DEBE ser revisada y firmada por un perito antes de uso legal."""


def imprimir_banner():
    print("""
    ========================================================
      ███████╗███████╗███████╗████████╗███████╗███╗   ███╗
      ██╔════╝██╔════╝██╔════╝╚══██╔══╝██╔════╝████╗ ████║
      █████╗  █████╗  █████╗     ██║   █████╗  ██╔████╔██║
      ██╔══╝  ██╔══╝  ██╔══╝     ██║   ██╔══╝  ██║╚██╔╝██║
      ██║     ███████╗███████╗   ██║   ███████╗██║ ╚═╝ ██║
      ╚═╝     ╚══════╝╚══════╝   ╚═╝   ╚══════╝╚═╝     ╚═╝
     FOREN-SYS: ASISTENTE DE INTELIGENCIA ARTIFICIAL (V3.1)
    * LLM Synthesis | Ollama Engine | Optimizado Raspberry Pi 5 *
    ========================================================
    """)


def comprobar_ollama():
    """Verifica si el servidor de Ollama está encendido y el modelo existe. Intenta auto-arrancar si es local y está caído."""
    print(f"[*] Conectando con el motor de IA en {OLLAMA_BASE_URL} (modelo: {MODELO_LLM})...")
    
    def _verificar_y_mostrar_modelos():
        try:
            modelos = requests.get(OLLAMA_BASE_URL + "/api/tags", timeout=10).json()
            nombres = [m['name'] for m in modelos.get('models', [])]
            if MODELO_LLM in nombres or any(MODELO_LLM.split(':')[0] in n for n in nombres):
                print(f"[+] Modelo '{MODELO_LLM}': DISPONIBLE")
            else:
                print(f"[!] Advertencia: Modelo '{MODELO_LLM}' no encontrado.")
                print(f"    Modelos disponibles: {', '.join(nombres) if nombres else 'Ninguno'}")
        except Exception:
            pass

    try:
        respuesta = requests.get(OLLAMA_BASE_URL + "/", timeout=3)
        if respuesta.status_code == 200:
            print("[+] Servidor Ollama: EN LÍNEA")
            _verificar_y_mostrar_modelos()
            return "en_linea"
    except requests.exceptions.ConnectionError:
        print(f"[-] ERROR: No se pudo conectar a Ollama en {OLLAMA_BASE_URL}.")
        
        # Intentar auto-arranque solo si es un motor local
        if "localhost" in OLLAMA_BASE_URL or "127.0.0.1" in OLLAMA_BASE_URL:
            print("[*] Motor de IA apagado. Intentando auto-arranque (Systemd/Binario)...")
            try:
                subprocess.run(["sudo", "-n", "systemctl", "start", "ollama"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=5)
            except Exception:
                pass
            try:
                subprocess.Popen(["ollama", "serve"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            except Exception:
                pass
                
            print("[*] Esperando a que el motor levante (max 15s)...")
            for _ in range(15):
                time.sleep(1)
                try:
                    if requests.get(OLLAMA_BASE_URL + "/", timeout=1).status_code == 200:
                        print("[+] Motor auto-arrancado con éxito.")
                        _verificar_y_mostrar_modelos()
                        return "auto"
                except Exception:
                    pass
            print("[-] ERROR CRÍTICO: Imposible arrancar Ollama. Verifica que el software esté instalado.")
        else:
            print("    • Motor local (RPi5): sudo systemctl start ollama")
            print("    • Motor remoto (PC):  Asegúrate de que Ollama está corriendo con OLLAMA_HOST=0.0.0.0")
        
        return False
    return False

def obtener_modelo_disponible():
    """Devuelve el modelo configurado si existe, o el primer modelo disponible como fallback."""
    try:
        modelos = requests.get(OLLAMA_BASE_URL + "/api/tags", timeout=5).json()
        nombres = [m['name'] for m in modelos.get('models', [])]
        if not nombres:
            return None
        if MODELO_LLM in nombres:
            return MODELO_LLM
        # Auto-selección del primero
        return nombres[0]
    except Exception:
        return None

def obtener_lista_modelos_ia():
    """Comprueba que Ollama esté online (auto-arrancando si es necesario) y devuelve la lista de modelos instalados."""
    estado = comprobar_ollama()
    if not estado:
        return []
    
    try:
        modelos = requests.get(OLLAMA_BASE_URL + "/api/tags", timeout=5).json()
        nombres = [m['name'] for m in modelos.get('models', [])]
        
        # Apagamos Ollama si lo acabamos de arrancar solo para listar los modelos
        # Así no dejamos la memoria bloqueada
        if estado == "auto":
            try:
                subprocess.run(["sudo", "-n", "systemctl", "stop", "ollama"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=5)
                subprocess.run(["pkill", "-f", "ollama serve"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=5)
            except Exception:
                pass
                
        return nombres
    except Exception:
        return []

def validar_modelo_ollama(modelo_nombre):
    """
    Valida que el modelo esté en la whitelist de modelos auditados. (CRÍTICA 5)
    Si no está, advierte pero permite continuar (el perito debe tomar la decisión).
    """
    if modelo_nombre in MODELOS_AUDITADOS:
        print(f"[+] Modelo '{modelo_nombre}' validado en whitelist de modelos auditados.")
        return True
    else:
        print(f"[!] ADVERTENCIA: Modelo '{modelo_nombre}' no está en la whitelist de modelos auditados.")
        print(f"    Modelos validados: {', '.join(sorted(MODELOS_AUDITADOS))}")
        print("    Continúa bajo responsabilidad del operador. Los resultados deben")
        print("    ser verificados especialmente antes de cualquier uso legal.")
        return False  # No bloquea, solo advierte


def sanitizar_prompt_injection(texto):
    """
    Previene inyección de prompts maliciosos en la evidencia. (CRÍTICA 6)
    Reemplaza patrones de jailbreak por marcadores inofensivos.
    """
    patrones_peligrosos = [
        (r'\[IGNORAR[^\]]*\]',         '[PATRÓN_BLOQUEADO]'),
        (r'\[IGNORE[^\]]*\]',          '[PATRÓN_BLOQUEADO]'),
        (r'(?i)ignore\s+previous',     '[PATRÓN_BLOQUEADO]'),
        (r'(?i)forget\s+previous',     '[PATRÓN_BLOQUEADO]'),
        (r'(?i)jailbreak',             '[PATRÓN_BLOQUEADO]'),
        (r'(?i)roleplay\s+as',         '[PATRÓN_BLOQUEADO]'),
        (r'(?i)act\s+as\s+if',        '[PATRÓN_BLOQUEADO]'),
        (r'(?i)--IgnoreInstructions',  '[PATRÓN_BLOQUEADO]'),
        (r'(?i)DAN\s+mode',            '[PATRÓN_BLOQUEADO]'),
    ]
    detecciones = []
    for patron, reemplazo in patrones_peligrosos:
        nuevas = re.findall(patron, texto)
        if nuevas:
            detecciones.extend(nuevas)
            texto = re.sub(patron, reemplazo, texto)
    if detecciones:
        print(f"[!] ALERTA: {len(detecciones)} patrones de prompt-injection detectados y neutralizados.")
        print(f"    Patrones: {detecciones[:5]}")
    return texto


def detectar_alucinaciones(respuesta_ia, evidencia_original):
    """
    Detecta posibles alucinaciones comparando nombres de archivo y años
    mencionados por la IA contra la evidencia original. (CRÍTICA 7)
    Retorna dict con métricas de confianza.
    """
    # Comparar años de 4 dígitos (2000-2030)
    años_evidencia = set(re.findall(r'\b(20[0-2]\d)\b', evidencia_original))
    años_ia        = set(re.findall(r'\b(20[0-2]\d)\b', respuesta_ia))
    años_nuevos    = años_ia - años_evidencia

    # Comparar nombres de archivo con extensión
    files_evidencia = set(f.lower() for f in re.findall(r'[a-zA-Z0-9_\-]{2,40}\.\w{2,5}', evidencia_original))
    files_ia        = set(f.lower() for f in re.findall(r'[a-zA-Z0-9_\-]{2,40}\.\w{2,5}', respuesta_ia))
    files_nuevos    = files_ia - files_evidencia - {'md', 'csv', 'txt', 'json', 'log'}

    total_anomalias = len(años_nuevos) + len(files_nuevos)

    if años_nuevos:
        print(f"[!] POSIBLE ALUCINACIÓN: IA mencionó años no en evidencia: {sorted(años_nuevos)}")
    if files_nuevos and len(files_nuevos) < 20:
        print(f"[!] POSIBLE ALUCINACIÓN: IA mencionó archivos no en evidencia: {sorted(files_nuevos)[:10]}")

    confianza = 'ALTA' if total_anomalias == 0 else ('MEDIA' if total_anomalias <= 3 else 'BAJA')
    if total_anomalias > 0:
        print(f"[!] Confianza en síntesis IA: {confianza} ({total_anomalias} anomalías detectadas)")
        print("    La síntesis REQUIERE revisión manual por el perito forense.")

    return {
        'años_nuevos':    sorted(años_nuevos),
        'archivos_nuevos': sorted(files_nuevos)[:20],
        'total_anomalias': total_anomalias,
        'confianza':       confianza,
    }


def _leer_csv_resumido(ruta, max_filas=20):
    """Lee un CSV y devuelve un resumen textual de sus primeras filas."""
    if not os.path.exists(ruta):
        return ""
    try:
        with open(ruta, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            header = next(reader, None)
            if not header:
                return ""
            filas = []
            for i, fila in enumerate(reader):
                if i >= max_filas:
                    break
                filas.append(fila)
        if not filas:
            return ""
        texto = " | ".join(header) + "\n"
        for fila in filas:
            texto += " | ".join(str(c)[:60] for c in fila) + "\n"
        return texto
    except Exception:
        return ""


def _agregar_bloque(evidencia, presupuesto, titulo, contenido):
    """Agrega un bloque de evidencia respetando el presupuesto de caracteres."""
    if not contenido or not contenido.strip():
        return evidencia, presupuesto
    bloque = f"--- {titulo} ---\n{contenido}--- FIN {titulo} ---\n\n"
    if len(bloque) < presupuesto:
        return evidencia + bloque, presupuesto - len(bloque)
    elif presupuesto > 200:
        bloque_truncado = bloque[:presupuesto - 50] + "\n[...TRUNCADO...]\n\n"
        return evidencia + bloque_truncado, 0
    return evidencia, presupuesto


def recopilar_inteligencia(carpeta_resultados):
    """Filtra y empaqueta TODA la evidencia V11.0 para la ventana de contexto de la IA."""
    print("[*] Recopilando evidencia V11.0 (9 fuentes de datos)...")

    evidencia_texto = ""
    presupuesto = MAX_CARACTERES_EVIDENCIA
    fuentes_cargadas = []

    # 1. REPORTE MAESTRO (Prioridad MÁXIMA)
    ruta_maestro = os.path.join(carpeta_resultados, "Reporte_Forense_Maestro.txt")
    if os.path.exists(ruta_maestro):
        with open(ruta_maestro, 'r', encoding='utf-8', errors='ignore') as f:
            contenido = f.read()
        idx = contenido.find("RESUMEN EJECUTIVO")
        fragmento = contenido[idx:][:1500] if idx != -1 else contenido[-1500:]
        idx_reg = contenido.find("[REGISTRY]")
        if idx_reg != -1:
            fragmento += "\n" + contenido[idx_reg:][:800]
        evidencia_texto, presupuesto = _agregar_bloque(evidencia_texto, presupuesto, "RESUMEN MAESTRO", fragmento)
        fuentes_cargadas.append("Reporte Maestro")

    # 2. USUARIOS DEL EQUIPO
    ruta_usuarios = os.path.join(carpeta_resultados, "Usuarios_Equipo.csv")
    texto_usuarios = _leer_csv_resumido(ruta_usuarios, 15)
    if texto_usuarios:
        evidencia_texto, presupuesto = _agregar_bloque(evidencia_texto, presupuesto, "USUARIOS DEL EQUIPO", texto_usuarios)
        fuentes_cargadas.append("Usuarios")

    # 3. HARDWARE Y USB
    ruta_hw = os.path.join(carpeta_resultados, "Hardware_y_USB.csv")
    texto_hw = _leer_csv_resumido(ruta_hw, 25)
    if texto_hw:
        evidencia_texto, presupuesto = _agregar_bloque(evidencia_texto, presupuesto, "HARDWARE Y DISPOSITIVOS USB", texto_hw)
        fuentes_cargadas.append("Hardware/USB")

    # 4. PERSISTENCIA
    ruta_persist = os.path.join(carpeta_resultados, "Mecanismos_Persistencia.csv")
    texto_persist = _leer_csv_resumido(ruta_persist, 15)
    if texto_persist:
        evidencia_texto, presupuesto = _agregar_bloque(evidencia_texto, presupuesto, "MECANISMOS DE PERSISTENCIA", texto_persist)
        fuentes_cargadas.append("Persistencia")

    # 5. ALERTAS DE SEGURIDAD del JSONL (Entropía alta)
    ruta_jsonl = os.path.join(carpeta_resultados, "Master_Timeline.jsonl")
    alertas = []
    if os.path.exists(ruta_jsonl):
        with open(ruta_jsonl, 'r', encoding='utf-8', errors='ignore') as f:
            for linea in f:
                try:
                    evento = json.loads(linea)
                    if evento.get("sospechoso_entropia") or evento.get("alerta_evasion"):
                        alertas.append(
                            f"- {evento.get('descripcion','?')} | "
                            f"Entropía: {evento.get('entropia', 'N/A')} | "
                            f"Ext. falsa: {'SÍ' if evento.get('alerta_evasion') else 'NO'}"
                        )
                except (json.JSONDecodeError, KeyError):
                    continue
    if alertas:
        evidencia_texto, presupuesto = _agregar_bloque(
            evidencia_texto, presupuesto, "ALERTAS SEGURIDAD (ENTROPÍA/EVASIÓN)",
            "\n".join(alertas[:20]) + "\n")
        fuentes_cargadas.append(f"Alertas ({len(alertas)})")

    # 6. PROGRAMAS INSTALADOS
    ruta_prog = os.path.join(carpeta_resultados, "Lista_Programas_Instalados.csv")
    texto_prog = _leer_csv_resumido(ruta_prog, 25)
    if texto_prog:
        evidencia_texto, presupuesto = _agregar_bloque(evidencia_texto, presupuesto, "PROGRAMAS INSTALADOS", texto_prog)
        fuentes_cargadas.append("Programas")

    # 7. HISTORIAL WEB
    ruta_web = os.path.join(carpeta_resultados, "Web_History_and_Bookmarks.txt")
    if os.path.exists(ruta_web) and presupuesto > 300:
        with open(ruta_web, 'r', encoding='utf-8', errors='ignore') as f:
            lineas_web = f.readlines()
        lineas_utiles = [l for l in lineas_web[2:40] if l.strip()]
        if lineas_utiles:
            evidencia_texto, presupuesto = _agregar_bloque(
                evidencia_texto, presupuesto, "HISTORIAL WEB",
                "".join(lineas_utiles))
            fuentes_cargadas.append("Web History")

    # 8. DESCARGAS WEB
    ruta_desc = os.path.join(carpeta_resultados, "Descargas_Web.csv")
    texto_desc = _leer_csv_resumido(ruta_desc, 15)
    if texto_desc and presupuesto > 200:
        evidencia_texto, presupuesto = _agregar_bloque(evidencia_texto, presupuesto, "DESCARGAS WEB", texto_desc)
        fuentes_cargadas.append("Descargas")

    # 9. EVENTOS DEL SISTEMA
    ruta_evt = os.path.join(carpeta_resultados, "Eventos_Sistema.csv")
    texto_evt = _leer_csv_resumido(ruta_evt, 15)
    if texto_evt and presupuesto > 200:
        evidencia_texto, presupuesto = _agregar_bloque(evidencia_texto, presupuesto, "EVENTOS SISTEMA (EVTX)", texto_evt)
        fuentes_cargadas.append("Event Logs")

    # 10. MULTIMEDIA CON METADATOS (GPS, cámaras)
    ruta_multi = os.path.join(carpeta_resultados, "Metadatos_Multimedia.csv")
    texto_multi = _leer_csv_resumido(ruta_multi, 10)
    if texto_multi and presupuesto > 200:
        evidencia_texto, presupuesto = _agregar_bloque(evidencia_texto, presupuesto, "MULTIMEDIA (EXIF/GPS)", texto_multi)
        fuentes_cargadas.append("Multimedia")

    # 11. ANÁLISIS VISUAL IA (si existe reporte de imágenes previo del mismo caso)
    reportes_vision = sorted(glob.glob(os.path.join(carpeta_resultados, 'Vision_IA_Completo_*.md')))
    if reportes_vision and presupuesto > 500:
        try:
            with open(reportes_vision[-1], 'r', encoding='utf-8', errors='ignore') as f:
                contenido_vision = f.read()
            # Incluir hasta 2000 chars del reporte visual más reciente
            evidencia_texto, presupuesto = _agregar_bloque(
                evidencia_texto, presupuesto, "ANÁLISIS VISUAL DE IMÁGENES (IA)", contenido_vision[:2000])
            fuentes_cargadas.append(f"Imágenes IA ({len(reportes_vision)} reportes)")
        except Exception as e:
            print(f"    [!] No se pudo leer reporte visual: {e}")

    caracteres_total = len(evidencia_texto)
    print(f"    [+] Fuentes cargadas: {', '.join(fuentes_cargadas)}")
    print(f"    [+] Evidencia empaquetada: {caracteres_total} caracteres (límite: {MAX_CARACTERES_EVIDENCIA})")

    return evidencia_texto


# ==========================================
# ANÁLISIS VISUAL MASIVO (VLM Multimodal)
# ==========================================

PROMPT_VISION_FORENSE = (
    "Eres un asistente forense digital. Analiza esta imagen recuperada de un "
    "dispositivo bajo investigación judicial. Describe con precisión:\n"
    "1) Todo texto legible visible en la imagen\n"
    "2) Objetos, personas o escenas identificables\n"
    "3) Fechas, horas, nombres o datos personales visibles\n"
    "4) Cualquier elemento que pueda ser relevante en una investigación\n\n"
    "REGLA CRÍTICA: No inventes ni asumas información que no esté visible "
    "claramente en la imagen. Si un elemento no es claro, indícalo explícitamente "
    "con la frase 'No distinguible'. Responde en español, máximo 200 palabras."
)

FORMATOS_IMAGEN = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp'}
TAMANO_MINIMO_BYTES = 20 * 1024   # 20 KB — descarta íconos/caché
MAX_IMAGENES_POR_CASO = 100
RUTAS_EXCLUIDAS = ('cache', 'thumb', '.thumbnails', 'temp', '__pycache__')


def _es_imagen_valida(ruta):
    """Devuelve True si la ruta es una imagen analizable (extensión + tamaño mínimo)."""
    if not os.path.isfile(ruta):
        return False
    if os.path.splitext(ruta)[1].lower() not in FORMATOS_IMAGEN:
        return False
    if os.path.getsize(ruta) < TAMANO_MINIMO_BYTES:
        return False
    ruta_lower = ruta.lower()
    if any(excl in ruta_lower for excl in RUTAS_EXCLUIDAS):
        return False
    return True


def analizar_imagenes_en_masa(carpeta_caso, ollama_url_base, modelo_vlm):
    """
    Escanea TODAS las imágenes forenses del caso, las analiza con el VLM
    remoto (gemma4:e4b u otro modelo multimodal) y guarda un reporte
    visual consolidado. Retorna un resumen de los hallazgos para incluir
    en la síntesis final.

    SOLO disponible con motor remoto. Si se detecta localhost, retorna ''
    con un aviso.

    Args:
        carpeta_caso: Ruta base del caso (contendrá subcarpetas con imágenes)
        ollama_url_base: URL del servidor Ollama (ej. http://192.168.2.3:11434)
        modelo_vlm: Nombre del modelo VLM a usar

    Returns:
        str: Texto con descripciones de todas las imágenes analizadas
    """
    # Seguridad: solo permitir en motor remoto
    if 'localhost' in ollama_url_base or '127.0.0.1' in ollama_url_base:
        print("[!] FASE VISUAL omitida: el análisis de imágenes requiere motor remoto (VLM no disponible en Raspberry Pi).")
        return ''

    print("[PROGRESO:35] FASE VISUAL: Escaneando imágenes del caso...")

    # Recopilar imágenes recursivamente
    imagenes_encontradas = []
    for raiz, _, archivos in os.walk(carpeta_caso):
        for archivo in archivos:
            ruta_completa = os.path.join(raiz, archivo)
            if _es_imagen_valida(ruta_completa):
                imagenes_encontradas.append(ruta_completa)

    if not imagenes_encontradas:
        print("    [!] No se encontraron imágenes válidas en el caso (mín. 20 KB). Fase visual omitida.")
        return ''

    # Ordenar por tamaño descendente (las más grandes son más relevantes)
    imagenes_encontradas.sort(key=lambda p: os.path.getsize(p), reverse=True)
    imagenes_a_analizar = imagenes_encontradas[:MAX_IMAGENES_POR_CASO]

    total = len(imagenes_a_analizar)
    omitidas = len(imagenes_encontradas) - total
    print(f"    [+] Imágenes encontradas: {len(imagenes_encontradas)} | A analizar: {total} | Omitidas por límite: {omitidas}")
    if omitidas > 0:
        print(f"    [!] Se omitieron {omitidas} imágenes menores por límite de seguridad ({MAX_IMAGENES_POR_CASO} max).")

    url_generate = ollama_url_base.rstrip('/') + "/api/generate"
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    ruta_reporte_vision = os.path.join(
        carpeta_caso,
        "03_Results_(Resultados_Extraidos)",
        f"Vision_IA_Completo_{ts}.md"
    )

    descripciones = []
    errores = 0

    for idx, ruta_img in enumerate(imagenes_a_analizar, 1):
        nombre = os.path.basename(ruta_img)
        pct = 35 + int((idx / total) * 20)  # Progreso de 35% a 55%
        print(f"[PROGRESO:{pct}] Analizando imagen {idx}/{total}: {nombre}")

        try:
            from PIL import Image
            import io
            with Image.open(ruta_img) as img:
                # Convertir a RGB por si tiene transparencia (PNG) o es CMYK/BMP
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                # Redimensionar si es extremadamente grande para no saturar VRAM
                img.thumbnail((1920, 1920))
                
                img_byte_arr = io.BytesIO()
                img.save(img_byte_arr, format='JPEG', quality=85)
                img_b64 = base64.b64encode(img_byte_arr.getvalue()).decode('utf-8')

            payload = {
                "model": modelo_vlm,
                "prompt": PROMPT_VISION_FORENSE,
                "images": [img_b64],
                "stream": False,
                "options": {
                    "temperature": 0.0,
                    "num_ctx": 98304,
                    "num_gpu": 42,
                    "top_k": 10,
                    "top_p": 0.5,
                    "repeat_penalty": 1.2
                }
            }

            resp = requests.post(url_generate, json=payload, timeout=120)
            resp.raise_for_status()
            descripcion = resp.json().get('response', '').strip()

            if descripcion:
                ruta_relativa = os.path.relpath(ruta_img, carpeta_caso)
                entrada = f"### Imagen: {nombre}\n**Ruta:** {ruta_relativa}\n\n{descripcion}\n"
                descripciones.append(entrada)
                print(f"    [+] Descripción obtenida ({len(descripcion)} chars)")
            else:
                print(f"    [!] Respuesta vacía para: {nombre}")
                errores += 1

        except requests.exceptions.Timeout:
            print(f"    [-] Timeout analizando: {nombre} (imagen muy compleja o red lenta)")
            errores += 1
        except Exception as e:
            print(f"    [-] Error en {nombre}: {e}")
            errores += 1

    # Guardar reporte visual consolidado
    if descripciones:
        cabecera = (
            f"# Reporte Visual IA — Análisis de Imágenes Forenses\n"
            f"**Caso:** {os.path.basename(carpeta_caso)}\n"
            f"**Fecha:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
            f"**Modelo VLM:** {modelo_vlm}\n"
            f"**Imágenes analizadas:** {len(descripciones)}/{total} "
            f"({'con errores: ' + str(errores) if errores else 'sin errores'})\n\n"
            f"---\n\n"
        )
        contenido_completo = cabecera + "\n".join(descripciones)
        try:
            os.makedirs(os.path.dirname(ruta_reporte_vision), exist_ok=True)
            with open(ruta_reporte_vision, 'w', encoding='utf-8') as f:
                f.write(contenido_completo)
            print(f"[PROGRESO:55] [+] Reporte visual guardado: {ruta_reporte_vision}")
        except Exception as e:
            print(f"    [-] No se pudo guardar el reporte visual: {e}")

        # Retornar resumen para incluir en prompt final (máx 3000 chars)
        resumen = "\n".join(descripciones)
        return resumen[:3000] if len(resumen) > 3000 else resumen
    else:
        print("    [-] No se obtuvo ninguna descripción visual válida.")
        return ''


def analizar_con_ia(evidencia_cruda, ruta_salida, ruta_auditoria, modelo_elegido=None):
    """
    Envía el prompt y evidencia al LLM con streaming.
    Genera:
      - Síntesis de inteligencia (ruta_salida) con DISCLAIMER visible
      - Registro de auditoría JSON firmado con SHA-256 (ruta_auditoria)  [CRÍTICA 4]
    Aplica sanitización anti-prompt-injection antes de construir el prompt.  [CRÍTICA 6]
    Detecta posibles alucinaciones al finalizar.                             [CRÍTICA 7]
    Retorna True si fue exitoso, False en caso contrario.
    """
    global MODELO_LLM
    
    estado_motor = comprobar_ollama()
    if not estado_motor:
        print("[-] Abortando síntesis porque el motor IA no está disponible o no pudo arrancar.")
        return False

    fue_autoarrancado = (estado_motor == "auto")

    # Override del modelo si se proporcionó uno válido, de lo contrario buscar automático
    if modelo_elegido:
        modelo_real = modelo_elegido
    else:
        modelo_real = obtener_modelo_disponible()
        
    if not modelo_real:
        print("[-] Abortando síntesis: No hay modelos instalados en Ollama.")
        return False
    
    MODELO_LLM = modelo_real

    # Sanitizar evidencia antes de inyectarla en el prompt
    evidencia_sanitizada = sanitizar_prompt_injection(evidencia_cruda)

    prompt_final = PROMPT_ASISTENTE.format(evidencia_cruda=evidencia_sanitizada)

    # Dynamic hardware capacity check
    is_local = "localhost" in OLLAMA_BASE_URL or "127.0.0.1" in OLLAMA_BASE_URL
    current_ctx = RPI_CTX if is_local else 98304
    current_threads = RPI_THREAD if is_local else PC_THREAD

    options = {
        "temperature": 0.0,
        "top_p": 0.5,
        "num_ctx": current_ctx,
        "num_thread": current_threads
    }
    
    if not is_local:
        options["num_gpu"] = 42
        options["top_k"] = 10
        options["repeat_penalty"] = 1.2

    carga_util = {
        "model": MODELO_LLM,
        "prompt": prompt_final,
        "stream": True,
        "options": options,
        "keep_alive": KEEP_ALIVE
    }

    # Registro de auditoría inicial [CRÍTICA 4]
    auditoria = {
        "timestamp_inicio_utc": datetime.utcnow().isoformat(),
        "modelo": MODELO_LLM,
        "motor_url": OLLAMA_BASE_URL,
        "parametros": options,
        "evidencia_sha256": hashlib.sha256(evidencia_cruda.encode('utf-8', errors='replace')).hexdigest(),
        "prompt_sha256":    hashlib.sha256(prompt_final.encode('utf-8', errors='replace')).hexdigest(),
        "tokens_count":     0,
        "sintesis_sha256":  None,
        "alucinaciones":    {},
        "timestamp_fin_utc": None,
        "estado": "EN_PROGRESO",
    }

    print("[*] Transmitiendo datos a la IA (streaming activado)...")
    print(f"[*] Motor: {OLLAMA_BASE_URL} | Modelo: {MODELO_LLM} | ctx: {NUM_CTX} | threads: {NUM_THREAD}")
    print("[*] La respuesta aparece en tiempo real. Ctrl+C para cancelar.\n")
    print("=" * 60)

    texto_completo = []

    try:
        respuesta = requests.post(OLLAMA_URL, json=carga_util, stream=True, timeout=TIMEOUT_SOLICITUD)
        respuesta.raise_for_status()

        for linea in respuesta.iter_lines():
            if linea:
                try:
                    fragmento = json.loads(linea)
                    token = fragmento.get("response", "")
                    texto_completo.append(token)
                    auditoria["tokens_count"] += 1
                    sys.stdout.write(token)
                    sys.stdout.flush()
                    if fragmento.get("done", False):
                        break
                except json.JSONDecodeError:
                    continue

        print("\n" + "=" * 60)

        sintesis_final = "".join(texto_completo)

        if sintesis_final.strip():
            # Detectar alucinaciones antes de guardar [CRÍTICA 7]
            resultado_alucinaciones = detectar_alucinaciones(sintesis_final, evidencia_cruda)
            auditoria["alucinaciones"] = resultado_alucinaciones

            # Construir cabecera del documento con DISCLAIMER visible [CRÍTICA 3]
            ts_gen = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            nivel_confianza = resultado_alucinaciones.get('confianza', 'MEDIA')
            cabecera = (
                f"{DISCLAIMER_LEGAL}\n\n"
                f"---\n"
                f"**Generado:** {ts_gen} | **Modelo:** {MODELO_LLM} | "
                f"**Confianza IA:** {nivel_confianza}\n\n"
                f"**ESTADO:** BORRADOR — REQUIERE REVISIÓN Y FIRMA DE PERITO CERTIFICADO\n\n"
                f"---\n\n"
                f"# SÍNTESIS DE INTELIGENCIA FORENSE\n\n"
                f"> ⚠️ Este documento es una síntesis automática. "
                f"No es un Dictamen Pericial. Requiere validación por perito.\n\n"
            )

            with open(ruta_salida, 'w', encoding='utf-8') as f:
                f.write(cabecera)
                f.write(sintesis_final)

            # Calcular hash de la síntesis y finalizar auditoría
            auditoria["sintesis_sha256"] = hashlib.sha256(
                sintesis_final.encode('utf-8', errors='replace')
            ).hexdigest()
            auditoria["timestamp_fin_utc"] = datetime.utcnow().isoformat()
            auditoria["estado"] = "COMPLETADO"

            print(f"\n[+] Síntesis completada. Guardada en:")
            print(f"    -> {ruta_salida}")
            print(f"[+] Confianza de síntesis: {nivel_confianza}")
            exito = True

        else:
            print("[-] La IA no generó contenido. Verifica que el modelo esté correctamente instalado.")
            auditoria["estado"] = "SIN_CONTENIDO"
            exito = False

    except requests.exceptions.Timeout:
        print(f"\n[-] TIMEOUT: La IA tardó más de {TIMEOUT_SOLICITUD // 60} minutos.")
        auditoria["estado"] = "TIMEOUT"
        exito = False
    except requests.exceptions.ConnectionError:
        print("\n[-] Se perdió la conexión con Ollama durante la generación.")
        print("    Posible causa: Ollama está apagado o se quedó sin memoria RAM (OOM Killer).")
        print("    Solución: Inicia Ollama o usa un modelo más pequeño.")
        auditoria["estado"] = "ERROR_CONEXION"
        exito = False
    except KeyboardInterrupt:
        print("\n\n[!] Generación cancelada por el usuario.")
        parcial = "".join(texto_completo)
        if parcial.strip():
            ruta_parcial = ruta_salida.replace(".md", "_PARCIAL.md")
            with open(ruta_parcial, 'w', encoding='utf-8') as f:
                f.write(f"{DISCLAIMER_LEGAL}\n\n<!-- INFORME PARCIAL - Cancelado por usuario -->\n\n{parcial}")
            print(f"    Síntesis parcial guardada en: {ruta_parcial}")
        auditoria["estado"] = "CANCELADO"
        exito = False
    except requests.exceptions.RequestException as e:
        print(f"[-] Error de comunicación con la IA: {e}")
        auditoria["estado"] = f"ERROR: {e}"
        exito = False
    finally:
        # Siempre guardar la auditoría, independientemente del resultado [CRÍTICA 4]
        try:
            with open(ruta_auditoria, 'w', encoding='utf-8') as f:
                json.dump(auditoria, f, indent=2, ensure_ascii=False)
            os.chmod(ruta_auditoria, 0o600)
            print(f"[+] Registro de auditoría guardado en:\n    -> {ruta_auditoria}")
        except Exception as e:
            print(f"[!] No se pudo guardar la auditoría: {e}")

        # Apagar Ollama si lo encendimos nosotros (para ahorrar RAM en la Raspberry Pi)
        if fue_autoarrancado:
            print("\n[*] Apagando el motor IA local para liberar memoria RAM...")
            try:
                subprocess.run(["sudo", "-n", "systemctl", "stop", "ollama"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=5)
                subprocess.run(["pkill", "-f", "ollama serve"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=5)
            except Exception as e:
                print(f"[!] No se pudo detener Ollama automáticamente: {e}")

    return exito

# ==========================================
# ANÁLISIS DOCUMENTAL INTELIGENTE (Document Intelligence)
# ==========================================

PROMPT_DOCS_FORENSE = (
    "Eres un perito informático forense. Analiza el siguiente texto extraído de un documento "
    "recuperado de una evidencia digital. Resume los puntos clave y extrae cualquier dato "
    "de posible interés judicial (nombres de personas, empresas, correos, teléfonos, fechas, transacciones "
    "o intenciones ocultas).\n\n"
    "REGLA CRÍTICA: Sé conciso y objetivo. No alucines información.\n\n"
    "Texto del documento:\n"
    "{texto_documento}"
)

def _extraer_texto_documento(ruta):
    """Intenta extraer texto de varios formatos ofimáticos y bases de datos."""
    ext = os.path.splitext(ruta)[1].lower()
    texto_extraido = ""
    try:
        if ext == '.pdf':
            import PyPDF2
            with open(ruta, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for i in range(min(len(reader.pages), 50)): # Max 50 paginas
                    t = reader.pages[i].extract_text()
                    if t: texto_extraido += t + "\n"
        elif ext == '.docx':
            import docx
            doc = docx.Document(ruta)
            texto_extraido = "\n".join([p.text for p in doc.paragraphs])
        elif ext == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(ruta, data_only=True, read_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_vals = [str(v) for v in row if v is not None]
                    if row_vals: texto_extraido += " | ".join(row_vals) + "\n"
        elif ext == '.pptx':
            import pptx
            prs = pptx.Presentation(ruta)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texto_extraido += shape.text + "\n"
        elif ext == '.msg':
            import extract_msg
            msg = extract_msg.Message(ruta)
            texto_extraido = f"De: {msg.sender}\nPara: {msg.to}\nFecha: {msg.date}\nAsunto: {msg.subject}\n\n{msg.body}"
        elif ext == '.eml':
            import email
            with open(ruta, 'rb') as f:
                msg = email.message_from_binary_file(f)
                texto_extraido = f"De: {msg.get('From')}\nPara: {msg.get('To')}\nAsunto: {msg.get('Subject')}\n\n"
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        texto_extraido += part.get_payload(decode=True).decode(errors='ignore')
        elif ext in ['.db', '.sqlite', '.sqlite3']:
            import sqlite3
            conn = sqlite3.connect(ruta)
            cursor = conn.cursor()
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
            tablas = cursor.fetchall()
            for t in tablas:
                tabla = t[0]
                texto_extraido += f"--- TABLA: {tabla} ---\n"
                cursor.execute(f"PRAGMA table_info({tabla});")
                cols = [c[1] for c in cursor.fetchall()]
                texto_extraido += "Columnas: " + ", ".join(cols) + "\n"
                cursor.execute(f"SELECT * FROM {tabla} LIMIT 20;")
                filas = cursor.fetchall()
                for fila in filas:
                    texto_extraido += str(fila) + "\n"
            conn.close()
        elif ext in ['.txt', '.csv', '.json', '.xml']:
            with open(ruta, 'r', encoding='utf-8', errors='ignore') as f:
                texto_extraido = f.read()
    except Exception as e:
        return f"[Error extrayendo texto: {e}]"
    
    return texto_extraido[:100000]  # Limite máximo de caracteres por documento

def analizar_documentos_en_masa(carpeta_caso, ollama_url_base, modelo_llm):
    """
    Escanea documentos en el caso, extrae texto y lo envía a la IA remota para su resumen.
    Aprovecha la ventana de contexto masiva de 65536 tokens.
    """
    if 'localhost' in ollama_url_base or '127.0.0.1' in ollama_url_base:
        print("[!] FASE DOCUMENTAL omitida: requiere motor remoto (PC Potente) para manejar contextos grandes.")
        return ''

    print("[PROGRESO:36] FASE DOCUMENTAL: Escaneando documentos ofimáticos, correos y BDs...")
    formatos_soportados = {'.pdf', '.docx', '.xlsx', '.pptx', '.msg', '.eml', '.db', '.sqlite', '.sqlite3', '.txt', '.csv'}
    rutas_excluidas = ('cache', 'temp', '__pycache__', 'system32')
    
    docs_encontrados = []
    for raiz, _, archivos in os.walk(carpeta_caso):
        for archivo in archivos:
            ruta_completa = os.path.join(raiz, archivo)
            ext = os.path.splitext(ruta_completa)[1].lower()
            if ext in formatos_soportados and os.path.getsize(ruta_completa) > 100:
                if not any(excl in ruta_completa.lower() for excl in rutas_excluidas):
                    docs_encontrados.append(ruta_completa)
    
    if not docs_encontrados:
        print("    [!] No se encontraron documentos analizables. Fase omitida.")
        return ''
        
    docs_encontrados.sort(key=lambda p: os.path.getsize(p), reverse=True)
    docs_a_analizar = docs_encontrados[:20] # Max 20 documentos más grandes
    
    total = len(docs_a_analizar)
    print(f"    [+] Documentos encontrados: {len(docs_encontrados)} | A analizar por IA: {total}")
    
    url_generate = ollama_url_base.rstrip('/') + "/api/generate"
    descripciones = []
    
    for idx, ruta_doc in enumerate(docs_a_analizar, 1):
        nombre = os.path.basename(ruta_doc)
        pct = 36 + int((idx / total) * 15) # Progreso de 36 a 51%
        print(f"[PROGRESO:{pct}] Analizando documento {idx}/{total}: {nombre}")
        
        texto_crudo = _extraer_texto_documento(ruta_doc)
        if not texto_crudo.strip():
            print("        [-] No se pudo extraer texto. Omitido.")
            continue
            
        payload = {
            "model": modelo_llm,
            "prompt": PROMPT_DOCS_FORENSE.format(texto_documento=texto_crudo),
            "stream": False,
            "options": {
                "temperature": 0.0,
                "num_ctx": 98304,
                "num_gpu": 42,
                "top_k": 10,
                "top_p": 0.5,
                "repeat_penalty": 1.2
            }
        }
        
        try:
            resp = requests.post(url_generate, json=payload, timeout=300)
            if resp.status_code == 200:
                res_json = resp.json()
                analisis = res_json.get("response", "").strip()
                descripciones.append(f"**Documento:** {nombre}\n**Ruta:** {ruta_doc.replace(carpeta_caso, '')}\n**Análisis Inteligente:**\n{analisis}\n")
                print(f"        [+] Resumen de {len(analisis)} caracteres generado.")
            else:
                print(f"        [-] Error de motor: {resp.status_code}")
        except Exception as e:
            print(f"        [-] Error de red/timeout: {e}")
            
    if not descripciones:
        return ''
        
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    ruta_reporte = os.path.join(carpeta_caso, "03_Results_(Resultados_Extraidos)", f"Docs_IA_Completo_{ts}.md")
    contenido_final = "# ANÁLISIS DOCUMENTAL INTELIGENTE (DOCUMENT INTELLIGENCE)\n\n" + "\n---\n".join(descripciones)
    
    try:
        with open(ruta_reporte, 'w', encoding='utf-8') as f:
            f.write(contenido_final)
        print(f"[PROGRESO:52] [+] Reporte documental guardado: {ruta_reporte}")
    except Exception as e:
        print(f"    [!] No se pudo guardar reporte documental: {e}")
        
    return contenido_final

# ==========================================
# INICIO — Modo no interactivo (para Web o CLI)
# ==========================================
if __name__ == "__main__":

    parser = argparse.ArgumentParser(description="Asistente IA Forense (Ollama) — Foren-Sys")
    parser.add_argument("--caso",    required=True,  help="ID del caso (Ej: TEST-123)")
    parser.add_argument("--dest",    required=False, default=DIRECTORIO_DEFAULT,
                        help=f"Ruta base donde vive el caso (default: {DIRECTORIO_DEFAULT})")
    parser.add_argument("--motor",   required=False, default="local",
                        choices=["local", "remoto"],
                        help="Motor Ollama: 'local' = RPi5, 'remoto' = PC Escritorio (default: local)")
    parser.add_argument("--host",    required=False, default=None,
                        help="URL del servidor Ollama remoto (Ej: http://192.168.1.50:11434). Sobreescribe --motor.")
    parser.add_argument("--model",   required=False, default=None,
                        help="Nombre del modelo a usar (Ej: gemma3:4b, llama3.2:3b, mistral:7b)")
    parser.add_argument("--ctx",     required=False, type=int, default=None,
                        help="Ventana de contexto en tokens (sobreescribe el perfil del motor)")
    parser.add_argument("--threads", required=False, type=int, default=None,
                        help="Número de hilos CPU (sobreescribe el perfil del motor)")
    parser.add_argument("--vision",  action="store_true", default=False,
                        help="Activar análisis visual masivo de imágenes (requiere motor remoto VLM)")
    parser.add_argument("--docs",  action="store_true", default=False,
                        help="Activar análisis inteligente de documentos, bases de datos y correos (extrae texto y resume)")
    args = parser.parse_args()

    # Validar y limpiar ID de caso (anti path-traversal)
    caso_id = re.sub(r'[^a-zA-Z0-9_\-]', '', args.caso.strip())
    if not caso_id or caso_id != args.caso.strip():
        print(f"[X] ID de caso inválido o contiene caracteres no permitidos: '{args.caso}'")
        print("    Permitido: alfanuméricos, guiones, guiones_bajos")
        sys.exit(1)

    directorio_base_actual = args.dest.strip()

    # Aplicar perfil de motor
    if args.host:
        OLLAMA_BASE_URL   = args.host.rstrip('/')
        OLLAMA_URL        = OLLAMA_BASE_URL + "/api/generate"
        NUM_CTX           = args.ctx     or PC_CTX
        NUM_THREAD        = args.threads or PC_THREAD
        TIMEOUT_SOLICITUD = PC_TIMEOUT
        perfil_nombre     = f"Personalizado ({OLLAMA_BASE_URL})"
    elif args.motor == "remoto":
        config_paths = [
            os.path.join(directorio_base_actual, ".ia_config.json"),
            os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "Casos_ForenSys", ".ia_config.json"),
            "/mnt/Destino_ForenSys/.ia_config.json"
        ]
        config_path = None
        for p in config_paths:
            if os.path.exists(p):
                config_path = p
                break
                
        if not config_path:
            print("[-] ERROR: Motor remoto seleccionado pero no hay host configurado.")
            print(f"    Crea un archivo .ia_config.json con la IP de tu PC de escritorio,")
            print("    o usa: --host http://192.168.X.X:11434")
            sys.exit(1)
        with open(config_path, 'r') as f:
            ia_cfg = json.load(f)
        host_remoto = ia_cfg.get('remote_host', '').rstrip('/')
        if not host_remoto:
            print("[-] ERROR: 'remote_host' está vacío en la configuración.")
            sys.exit(1)
        OLLAMA_BASE_URL   = host_remoto
        OLLAMA_URL        = OLLAMA_BASE_URL + "/api/generate"
        NUM_CTX           = args.ctx     or ia_cfg.get('ctx',     PC_CTX)
        NUM_THREAD        = args.threads or ia_cfg.get('threads', PC_THREAD)
        TIMEOUT_SOLICITUD = ia_cfg.get('timeout', PC_TIMEOUT)
        perfil_nombre     = f"PC Escritorio Remoto ({OLLAMA_BASE_URL})"
    else:
        OLLAMA_BASE_URL   = "http://localhost:11434"
        OLLAMA_URL        = OLLAMA_BASE_URL + "/api/generate"
        NUM_CTX           = args.ctx     or RPI_CTX
        NUM_THREAD        = args.threads or RPI_THREAD
        TIMEOUT_SOLICITUD = RPI_TIMEOUT
        perfil_nombre     = "Raspberry Pi 5 (Local)"

    if args.model:
        MODELO_LLM = args.model.strip()

    imprimir_banner()
    print(f"[PROGRESO:5] Iniciando Análisis Asistido por IA para el caso: {caso_id}")
    print(f"    Motor:   {perfil_nombre}")
    print(f"    Modelo:  {MODELO_LLM}")
    print(f"    Config:  ctx={NUM_CTX} | threads={NUM_THREAD} | timeout={TIMEOUT_SOLICITUD}s")
    print(f"    Evidencia máx.: {MAX_CARACTERES_EVIDENCIA} caracteres")

    # Validar modelo contra whitelist [CRÍTICA 5]
    print(f"[PROGRESO:8] Validando modelo contra whitelist de modelos auditados...")
    validar_modelo_ollama(MODELO_LLM)  # Advierte pero no bloquea

    # Verificar Ollama
    print(f"[PROGRESO:10] Verificando motor Ollama...")
    if not comprobar_ollama():
        print("[-] ERROR: Motor Ollama no disponible.")
        sys.exit(1)

    # Localizar carpeta de resultados
    carpeta_resultados = os.path.join(directorio_base_actual, caso_id, "03_Results_(Resultados_Extraidos)")
    if not os.path.exists(carpeta_resultados):
        print(f"[-] ERROR: No se encontró la carpeta de resultados en: {carpeta_resultados}")
        print("    Asegúrate de haber ejecutado el Módulo 7 (Normalización) primero.")
        sys.exit(1)

    # Rutas de salida: síntesis + auditoría [CRÍTICA 4]
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    ruta_sintesis  = os.path.join(carpeta_resultados, f"Sintesis_IA_{caso_id}.md")
    ruta_auditoria = os.path.join(carpeta_resultados, f"Auditoria_IA_{caso_id}_{ts}.json")

    # Recopilar evidencia
    print("[PROGRESO:30] Recopilando y filtrando evidencia de los módulos anteriores...")
    evidencia_filtrada = recopilar_inteligencia(carpeta_resultados)
    if len(evidencia_filtrada) < 50:
        print("[-] Advertencia: Poca evidencia disponible. ¿Ejecutaste el Módulo 7?")
        print("[-] Continuando con la evidencia disponible...")

    # FASE 2 (OPCIONAL): Análisis Visual Masivo
    resumen_visual = ''
    if args.vision:
        carpeta_caso_base = os.path.join(directorio_base_actual, caso_id)
        resumen_visual = analizar_imagenes_en_masa(carpeta_caso_base, OLLAMA_BASE_URL, MODELO_LLM)
        if resumen_visual:
            # Inyectar el resumen visual en la evidencia textual antes de la síntesis
            bloque_vision = (
                "\n--- ANÁLISIS VISUAL DE IMÁGENES (IA) ---\n"
                + resumen_visual
                + "\n--- FIN ANÁLISIS VISUAL ---\n"
            )
            evidencia_filtrada += bloque_vision
            print(f"[+] Análisis visual integrado en la evidencia ({len(resumen_visual)} chars adicionales).")
    else:
        print("[INFO] Análisis visual omitido (usa --vision para activarlo con motor remoto).")

    # FASE 2.5 (OPCIONAL): Análisis Documental Masivo
    resumen_documental = ''
    if args.docs:
        carpeta_caso_base = os.path.join(directorio_base_actual, caso_id)
        resumen_documental = analizar_documentos_en_masa(carpeta_caso_base, OLLAMA_BASE_URL, MODELO_LLM)
        if resumen_documental:
            # Inyectar el resumen documental en la evidencia textual antes de la síntesis
            bloque_docs = (
                "\n--- ANÁLISIS DE DOCUMENTOS (IA) ---\n"
                + resumen_documental
                + "\n--- FIN ANÁLISIS DE DOCUMENTOS ---\n"
            )
            evidencia_filtrada += bloque_docs
            print(f"[+] Análisis documental integrado en la evidencia ({len(resumen_documental)} chars adicionales).")
    else:
        print("[INFO] Análisis documental omitido (usa --docs para activarlo con motor remoto).")

    # Análisis IA (FASE 3: Síntesis final)
    print(f"[PROGRESO:60] Transmitiendo evidencia al LLM (motor: {perfil_nombre})...")
    analizar_con_ia(evidencia_filtrada, ruta_sintesis, ruta_auditoria)

    if os.path.exists(ruta_sintesis):
        print(f"[PROGRESO:100] Síntesis de inteligencia guardada en: {ruta_sintesis}")
        print(f"               Auditoría guardada en: {ruta_auditoria}")
        print()
        print("  ⚠️  RECUERDE: Esta síntesis requiere revisión y firma de perito certificado")
        print("     antes de ser usada en procedimientos legales.")
    else:
        print("[-] ERROR: No se generó la síntesis final.")
        sys.exit(1)
